from fastapi.responses import FileResponse
from python_pptx_text_replacer import TextReplacer
from pptx.chart.data import CategoryChartData
from pptx import Presentation
from fastapi import FastAPI

app = FastAPI()

# Function to replace chart data on a specific slide
def replace_chart_with_data(slide, chart_index, chart_data):
    chart_count = 0
    for shape in slide.shapes:
        if shape.has_chart:
            chart_count += 1
            if chart_count == chart_index + 1:  
                chart = shape.chart
                chart.replace_data(chart_data)
                print(f"Chart with index {chart_index} found and replaced successfully on the specified slide.")
                return
    print(f"Chart with index {chart_index} not found on the specified slide.")

# Endpoint to generate the presentation
@app.post("/generate")
async def generate(data: dict):
    # Open the PowerPoint presentation to replace charts
    prs = Presentation('Cobatext3.pptx')

    # Replace chart data for top 10 online media
    top_10_online_media = data['result']['top_10_online_media']
    sorted_online_media = sorted(top_10_online_media.items(), key=lambda x: x[1], reverse=False)
    categories = [item[0] for item in sorted_online_media]
    values = [item[1] for item in sorted_online_media]
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('', values)
    slide_index = 2
    chart_index_to_replace = 0
    for i, slide in enumerate(prs.slides):
        if i == slide_index:
            replace_chart_with_data(slide, chart_index_to_replace, chart_data)
            break

    # Replace chart data for top 10 printed media
    top_10_printed_media = data['result']['top_10_printed_media']
    sorted_printed_media = sorted(top_10_printed_media.items(), key=lambda x: x[1], reverse=False)
    categories = [item[0] for item in sorted_printed_media]
    values = [item[1] for item in sorted_printed_media]
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('', values)
    slide_index = 2
    chart_index_to_replace = 1
    for i, slide in enumerate(prs.slides):
        if i == slide_index:
            replace_chart_with_data(slide, chart_index_to_replace, chart_data)
            break

    # Replace chart data for per day detail
    per_day_detail = data['result']['per_day_detail']
    dates = []
    online_values = []
    printed_values = []
    for date, details in per_day_detail.items():
        dates.append(date)
        online_values.append(details['online'])
        printed_values.append(details['printed'])
    chart_data = CategoryChartData()
    chart_data.categories = dates
    chart_data.add_series('', online_values)
    chart_data.add_series('', printed_values)
    slide_index = 2
    chart_index_to_replace = 2
    for i, slide in enumerate(prs.slides):
        if i == slide_index:
            replace_chart_with_data(slide, chart_index_to_replace, chart_data)
            break

    # Replace chart data for top 10 online influencer
    top_10_online_influencer = data['result']['top_10_online_influencer']
    sorted_online_influencer = sorted(top_10_online_influencer.items(), key=lambda x: x[1], reverse=False)
    categories = [item[0] for item in sorted_online_influencer]
    values = [item[1] for item in sorted_online_influencer]
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('', values)
    slide_index = 4
    chart_index_to_replace = 0
    for i, slide in enumerate(prs.slides):
        if i == slide_index:
            replace_chart_with_data(slide, chart_index_to_replace, chart_data)
            break

    # Replace chart data for all days detail
    all_days_detail = data['result']['all_days_detail']
    categories = [item['text'] for item in all_days_detail]
    totals = [item['percentage'] for item in all_days_detail]
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('', totals)
    slide_index = 3
    chart_index_to_replace = 0
    for i, slide in enumerate(prs.slides):
        if i == slide_index:
            replace_chart_with_data(slide, chart_index_to_replace, chart_data)
            break

    # Replace chart data for sentiment
    sentiment = data['result']['sentiment']
    categories = ["Negatif" ,"Netral" ,"Positif"]
    totals = [sentiment['negative']['percentage'], sentiment['neutral']['percentage'], sentiment['positive']['percentage']]
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('', totals)
    slide_index = 3
    chart_index_to_replace = 1
    for i, slide in enumerate(prs.slides):
        if i == slide_index:
            replace_chart_with_data(slide, chart_index_to_replace, chart_data)
            break

    # Save the modified presentation
    save_file = 'template/template.pptx'
    prs.save(save_file)

    # Replace text in the presentation
    replacer = TextReplacer(save_file, slides='', tables=False, charts=False, textframes=True)
    replacer.replace_text([
        ('2.251', str(data['result']['total_online_news'])),
        ('761', str(data['result']['total_online_media'])),
        ('135', str(data['result']['total_printed_news'])),
        ('60', str(data['result']['total_printed_media'])),
        ('1 – 7 April 2022' , str(data['result']['earliest_date'] + " Sampai " + str(data['result']['latest_date']))),
        # Add more text replacements as needed
    ])

    # Save the modified presentation with replaced text
    file_output = "result/Pitching_Report_DUPK_BI_Pengaduan_Sistem_Pembayaran_1_7_April_2022.pptx"
    replacer.write_presentation_to_file(file_output)

    # Return the file for download
    return FileResponse(file_output, filename="Pitching_Report_DUPK_BI_Pengaduan_Sistem_Pembayaran_1_7_April_2022.pptx")
